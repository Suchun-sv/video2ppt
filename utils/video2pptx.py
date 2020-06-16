import cv2
import tqdm
from PIL import Image
import os
import sys
import shutil
from pptx import Presentation
from pptx.util import Inches

class video2pptx:
    def __init__(self,videoPath=None,pictureFolder=None,pptxPath=None,time_interval=None,reducePictureFolder=None):
        self.videoPath = videoPath
        os.chdir("..//../")
        self.pictureFolder = pictureFolder
        self.pptxPath = pptxPath
        self.time_interval = time_interval
        self.reducePictureFolder = reducePictureFolder
        
    def capFrame(self,pictureFolder=None,time_interval=None,videoPath=None,):
        # 视频路径的判断
        if videoPath == None and self.videoPath == None:
            raise FileNotFoundError("videoPath didn't specify")
        elif videoPath:
            self.videoPath = videoPath
        elif self.videoPath:
            videoPath = self.videoPath
        if not os.path.exists(videoPath):
            raise FileNotFoundError("cannot find {}".format(videoPath))
        
        #输出截图路径的判断,默认和视频路径一样
        if pictureFolder == None and self.pictureFolder == None:
            pictureFolder = videoPath
            self.pictureFolder = videoPath
        elif pictureFolder:
            self.pictureFolder = pictureFolder
        elif self.pictureFolder:
            pictureFolder = self.pictureFolder
        if not os.path.exists(pictureFolder):
            os.mkdir(pictureFolder)
        print("{} selected".format(pictureFolder))
        self.deleteFolder(pictureFolder)
        
        #视频抓取间隔时间的判断,默认为30秒
        if self.time_interval == None and time_interval == None:
            time_interval = 30
        elif self.time_interval:
            time_interval = self.time_interval
        
        cap = cv2.VideoCapture(videoPath)
        suc = cap.isOpened()  # 是否成功打开
        frame_count = 0
        frame_rate = cap.get(5)
        total_frames = int(cap.get(7))#总帧数
        time_interval_frames = int(time_interval*frame_rate)
#         for it in tqdm.tnrange(total_frames//time_interval_frames):
#             while True:
#                 frame_count += 1
#                 suc, frame = cap.read()
#                 params = []
#                 params.append(2)  # params.append(1)
#                 if frame_count%time_interval_frames == 0:
#                     cv2.imwrite('{}\\{:06}.jpg'.format(pictureFolder,int(frame_count//frame_rate)), frame, params)
#                     break
                    
        for it in tqdm.trange(total_frames//time_interval_frames):
            cap.set(cv2.CAP_PROP_POS_FRAMES ,frame_count)
            suc, frame = cap.read()
            params = []
            params.append(2)  # params.append(1)
            cv2.imwrite('{}\\{:06}.jpg'.format(pictureFolder,int(frame_count//frame_rate)), frame, params)
            frame_count += time_interval_frames
            
                    
        cap.release()
        print({"pictures has been send to {}".format(pictureFolder)})

    def hash_img(self,img):#计算图片的特征序列
        a=[]#存储图片的像素
        hash_img=''#特征序列
        width,height=10,10#图片缩放大小
        img=img.resize((width,height))#图片缩放为width×height
        for y in range(img.height):
            b=[]
            for x in range(img.width):
                pos=x,y
                color_array = img.getpixel(pos)#获得像素
                color=sum(color_array)/3#灰度化
                b.append(int(color))
            a.append(b)
        for y in range(img.height):
            avg=sum(a[y])/len(a[y])#计算每一行的像素平均值
            for x in range(img.width):
                if a[y][x]>=avg:#生成特征序列,如果此点像素大于平均值则为1,反之为0
                    hash_img+='1'
                else:
                    hash_img+='0'

        return hash_img

    def similar(self,img1,img2):#求相似度
        hash1=self.hash_img(img1)#计算img1的特征序列
        hash2=self.hash_img(img2)#计算img2的特征序列
        differnce=0
        for i in range(len(hash1)):
            differnce+=abs(int(hash1[i])-int(hash2[i]))
        similar=1-(differnce/len(hash1))
        return similar
    
    def calcSimilar(self,pictureFolder=None):
        print("calculating similarity....")
        if pictureFolder==None and self.pictureFolder==None:
            raise FileNotFoundError("pictureFolder didn't sepecify")
        elif self.pictureFolder:
            pictureFolder = self.pictureFolder
            
        # 计算相似度
        similar_score = [0]
        files = os.listdir(pictureFolder)
        for i in range(0,len(files)-1):
            img1=Image.open(os.path.join(pictureFolder,files[i]))
            img2=Image.open(os.path.join(pictureFolder,files[i+1]))
            similar_score.append(self.similar(img1,img2))
        return similar_score
    
    def copyPictureBySimilar(self,threshold,pictureFolder=None,reducePictureFolder=None,similar_score=None):
        
        #原始截图文件夹
        if pictureFolder:
            self.pictureFolder = pictureFolder
            
        if not similar_score:
            similar_score = self.calcSimilar(self.pictureFolder)       
        #print(similar_score)
        #精简截图文件夹路径的判断,默认在截图文件夹后加上reduce标识
        if reducePictureFolder == None and self.reducePictureFolder == None:
            reducePictureFolder = self.pictureFolder+"_reduce"
            self.reducePictureFolder = reducePictureFolder
        elif reducePictureFolder:
            pass
        elif self.reducePictureFolder:
            reducePictureFolder = self.reducePictureFolder
        if not os.path.exists(reducePictureFolder):
            os.mkdir(reducePictureFolder)
        self.deleteFolder(reducePictureFolder)
        
        files = os.listdir(pictureFolder)
        #print(files)
        for i in range(len(files)):
            if similar_score[i]<threshold:
                #print(os.path.join(self.pictureFolder,files[i])+os.path.join(reducePictureFolder,files[i]))
                shutil.copyfile(os.path.join(self.pictureFolder,files[i]),os.path.join(reducePictureFolder,files[i])) 
                
    def createPPtx(self,pptName,pptxTemplate='..//template//core.pptx',pictureFile=None):
        if os.path.exists(pptxTemplate):
            prs = Presentation(pptxTemplate)
        else:
            prs = Presentation()
    
        if pictureFile ==None and self.reducePictureFolder==None:
            raise FileNotFoundError("reducePictureFolder cannot find")
        elif self.reducePictureFolder:
            pictureFile = self.reducePictureFolder
        
        for i in os.listdir(pictureFile):
            slide= prs.slides.add_slide(prs.slide_layouts[6])
            img_path=os.path.join(pictureFile,i)
            # 文件路径
            left,top,width,height=Inches(0),Inches(0),Inches(11.02362205),Inches(5.90551181)
            # 预设位置及大小
            pic=slide.shapes.add_picture(img_path,left,top,width,height)
        if pptName.endswith(".pptx"):
            prs.save(pptName)
        else:
            prs.save(pptName+".pptx")
        print("{}  creaed!".format(pptName))
    
    def deleteFolder(self, folder):
        files = os.listdir(folder)
        files = [os.path.join(folder,x) for x in files]
        for file in files:
            os.remove(file)